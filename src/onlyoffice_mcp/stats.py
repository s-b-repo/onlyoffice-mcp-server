"""Document statistics — polymorphic per format."""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from . import cursor as cursor_mod

log = logging.getLogger(__name__)


_SENTENCE_RE = re.compile(r"[.!?]+(?:\s+|$)")
_WORD_RE = re.compile(r"\S+")


def _docx_stats(path: Path) -> dict[str, Any]:
    from docx import Document

    doc = Document(str(path))
    word_count = 0
    char_count = 0
    char_count_no_spaces = 0
    sentence_count = 0
    heading_count = 0
    table_count = len(doc.tables)
    image_count = 0
    hyperlink_count = 0
    comment_count = 0

    for para in doc.paragraphs:
        text = para.text
        word_count += len(_WORD_RE.findall(text))
        char_count += len(text)
        char_count_no_spaces += len(text.replace(" ", "").replace("\t", ""))
        sentence_count += len(_SENTENCE_RE.findall(text))
        if para.style and para.style.name.lower().startswith("heading"):
            heading_count += 1

    # Walk parts for images / hyperlinks / comments via lxml.
    try:
        from lxml import etree

        _wns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        body_xml = doc.element.body
        image_count = len(body_xml.findall(f".//{{{_wns}}}drawing"))
        hyperlink_count = len(body_xml.findall(f".//{{{_wns}}}hyperlink"))
        # Comments live in word/comments.xml — check if the part exists.
        for part in doc.part.package.parts:
            if part.partname and str(part.partname).endswith("/comments.xml"):
                from .safety import safe_parse_xml
                root = safe_parse_xml(part.blob)
                comment_count = len(root.findall(f".//{{{_wns}}}comment"))
                break
    except Exception as exc:
        log.warning("XML stats extraction failed for %s: %s", path, exc)

    return {
        "format": "docx",
        "word_count": word_count,
        "char_count": char_count,
        "char_count_no_spaces": char_count_no_spaces,
        "paragraph_count": len(doc.paragraphs),
        "sentence_count": sentence_count,
        "heading_count": heading_count,
        "page_count_estimate": cursor_mod.estimate_docx_pages(path),
        "table_count": table_count,
        "image_count": image_count,
        "hyperlink_count": hyperlink_count,
        "comment_count": comment_count,
    }


def _xlsx_stats(path: Path) -> dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=False)
    sheet_names = wb.sheetnames
    total_rows = 0
    total_cells_used = 0
    formula_count = 0
    chart_count = 0
    table_count = 0

    for name in sheet_names:
        ws = wb[name]
        # max_row gives the highest used row.
        total_rows += ws.max_row or 0
        chart_count += len(getattr(ws, "_charts", []) or [])
        table_count += len(getattr(ws, "tables", {}) or {})
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                total_cells_used += 1
                if cell.data_type == "f" or (
                    isinstance(cell.value, str) and cell.value.startswith("=")
                ):
                    formula_count += 1

    return {
        "format": "xlsx",
        "sheet_count": len(sheet_names),
        "sheet_names": sheet_names,
        "total_rows": total_rows,
        "total_cells_used": total_cells_used,
        "formula_count": formula_count,
        "chart_count": chart_count,
        "table_count": table_count,
    }


def _pptx_stats(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    total_text_chars = 0
    total_text_words = 0
    image_count = 0
    chart_count = 0
    notes_chars_total = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
            if shape.has_chart:
                chart_count += 1
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text or ""
                    total_text_chars += len(text)
                    total_text_words += len(_WORD_RE.findall(text))
        if slide.has_notes_slide:
            notes_chars_total += len(slide.notes_slide.notes_text_frame.text or "")

    return {
        "format": "pptx",
        "slide_count": len(prs.slides),
        "total_text_chars": total_text_chars,
        "total_text_words": total_text_words,
        "image_count": image_count,
        "chart_count": chart_count,
        "notes_chars_total": notes_chars_total,
    }


def stats(path: str) -> dict[str, Any]:
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(p)
    ext = p.suffix.lstrip(".").lower()
    if ext == "docx":
        return _docx_stats(p)
    if ext == "xlsx":
        return _xlsx_stats(p)
    if ext == "pptx":
        return _pptx_stats(p)
    raise ValueError(
        f"Unsupported format for stats: '.{ext}'.\n"
        f"Supported formats: docx, xlsx, pptx."
    )
