"""Page background, watermark, slide background and sheet tab colour.

DOCX page background and watermark require direct XML manipulation
(python-docx doesn't expose either). Slide background and sheet tab colour
have native APIs.

References:
- OOXML §17.2.1 — w:background element on w:document
- OOXML §17.15.1.20 — w:displayBackgroundShape in settings.xml
- OOXML §20.4.2.1 — wp:anchor element for behind-text images
- OOXML §20.1.8.1 — a:alphaModFix for image opacity
"""

from __future__ import annotations

import logging
import xml.sax.saxutils
from pathlib import Path

from lxml import etree

from .validation import validate_path, validate_color
from .safety import safe_parse_xml, ALLOWED_IMAGE_EXTENSIONS

log = logging.getLogger(__name__)

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
V_NS = "urn:schemas-microsoft-com:vml"
O_NS = "urn:schemas-microsoft-com:office:office"
WP_NS = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PIC_NS = "http://schemas.openxmlformats.org/drawingml/2006/picture"


def _normalise_color(color_hex: str) -> str:
    return validate_color(color_hex)


# --------------------------------------------------------------------------
# DOCX page background
# --------------------------------------------------------------------------

def docx_set_background(path: str, color_hex: str) -> str:
    """Set the page background colour on a .docx file.

    Inserts ``<w:background w:color="RRGGBB"/>`` as the first child of
    ``<w:document>`` AND ``<w:displayBackgroundShape/>`` in settings.xml so
    LibreOffice and Word Online actually render the colour.
    """
    from docx import Document

    color = _normalise_color(color_hex)
    p = validate_path(path, must_exist=True, expected_ext="docx", operation="set_background")

    doc = Document(str(p))
    root = doc.element  # <w:document>
    nsmap = {"w": W_NS}

    # Remove any existing w:background then prepend a new one.
    for existing in root.findall(f"{{{W_NS}}}background"):
        root.remove(existing)
    bg = etree.SubElement(root, f"{{{W_NS}}}background")
    bg.set(f"{{{W_NS}}}color", color)
    # SubElement appends; we need it as the first child.
    root.remove(bg)
    root.insert(0, bg)

    # Settings part: add w:displayBackgroundShape so the colour is shown.
    try:
        settings_part = doc.part.package.part_related_by_reltype(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/settings"
        )
    except Exception as exc:
        logging.getLogger(__name__).warning("settings part lookup failed: %s", exc)
        settings_part = None
    if settings_part is not None:
        settings_root = safe_parse_xml(settings_part.blob)
        if settings_root.find(f"{{{W_NS}}}displayBackgroundShape") is None:
            display = etree.SubElement(
                settings_root, f"{{{W_NS}}}displayBackgroundShape"
            )
            # Re-insert as an early child for canonical order.
            settings_root.remove(display)
            settings_root.insert(0, display)
            settings_part._blob = etree.tostring(  # noqa: no public setter
                settings_root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )

    doc.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# DOCX page background image (behind-text anchor in section headers)
# --------------------------------------------------------------------------

def _build_bg_anchor_xml(
    *,
    r_id: str,
    cx: int,
    cy: int,
    offset_x: int,
    offset_y: int,
    doc_pr_id: int,
    opacity_pct: int,
) -> str:
    """Build the OOXML wp:anchor XML for a behind-text background image."""
    if opacity_pct < 100:
        alpha_amt = opacity_pct * 1000
        alpha_xml = f'<a:alphaModFix xmlns:a="{A_NS}" amt="{alpha_amt}"/>'
    else:
        alpha_xml = ""

    return f"""\
<w:r xmlns:w="{W_NS}" xmlns:wp="{WP_NS}" xmlns:a="{A_NS}"
     xmlns:pic="{PIC_NS}" xmlns:r="{R_NS}">
  <w:rPr><w:noProof/></w:rPr>
  <w:drawing>
    <wp:anchor distT="0" distB="0" distL="0" distR="0"
               simplePos="0" relativeHeight="0" behindDoc="1"
               locked="1" layoutInCell="1" allowOverlap="1">
      <wp:simplePos x="0" y="0"/>
      <wp:positionH relativeFrom="page">
        <wp:posOffset>{offset_x}</wp:posOffset>
      </wp:positionH>
      <wp:positionV relativeFrom="page">
        <wp:posOffset>{offset_y}</wp:posOffset>
      </wp:positionV>
      <wp:extent cx="{cx}" cy="{cy}"/>
      <wp:effectExtent l="0" t="0" r="0" b="0"/>
      <wp:wrapNone/>
      <wp:docPr id="{doc_pr_id}" name="PageBackground"
                descr="Page background image"/>
      <wp:cNvGraphicFramePr>
        <a:graphicFrameLocks noChangeAspect="1"/>
      </wp:cNvGraphicFramePr>
      <a:graphic>
        <a:graphicData uri="{PIC_NS}">
          <pic:pic>
            <pic:nvPicPr>
              <pic:cNvPr id="{doc_pr_id}" name="PageBackground"/>
              <pic:cNvPicPr/>
            </pic:nvPicPr>
            <pic:blipFill>
              <a:blip r:embed="{r_id}">{alpha_xml}</a:blip>
              <a:stretch><a:fillRect/></a:stretch>
            </pic:blipFill>
            <pic:spPr>
              <a:xfrm>
                <a:off x="0" y="0"/>
                <a:ext cx="{cx}" cy="{cy}"/>
              </a:xfrm>
              <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
            </pic:spPr>
          </pic:pic>
        </a:graphicData>
      </a:graphic>
    </wp:anchor>
  </w:drawing>
</w:r>"""


def _remove_bg_images_from_header(header_element: etree._Element) -> None:
    """Remove any existing PageBackground anchors from a header element."""
    for p_el in list(header_element.findall(f"{{{W_NS}}}p")):
        for r_el in p_el.findall(f"{{{W_NS}}}r"):
            for drawing in r_el.findall(f"{{{W_NS}}}drawing"):
                for anchor in drawing.findall(f"{{{WP_NS}}}anchor"):
                    doc_pr = anchor.find(f"{{{WP_NS}}}docPr")
                    if doc_pr is not None and doc_pr.get("name") == "PageBackground":
                        header_element.remove(p_el)
                        return


def docx_set_background_image(
    path: str,
    image_path: str,
    *,
    page_size: str = "a4",
    landscape: bool = False,
    offset_x_mm: float = 0,
    offset_y_mm: float = 0,
    width_mm: float | None = None,
    height_mm: float | None = None,
    opacity: int = 100,
) -> str:
    """Set a background image on a .docx file.

    The image is placed behind all text via section headers. Calling this
    again replaces the previous background image (no stacking).

    Args:
        offset_x_mm / offset_y_mm: position offset from top-left in mm.
        width_mm / height_mm: image size in mm (defaults to page size).
        opacity: image opacity 1-100 (100 = fully opaque, 30 = very faded).
    """
    from docx import Document
    from docx.shared import Mm
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    from . import safety

    p = validate_path(
        path, must_exist=True, expected_ext="docx",
        operation="set_background_image",
    )
    img = Path(image_path).expanduser().resolve()
    safety.check_path_safety(img)
    if not img.exists():
        raise ValueError(
            f"Image not found: {img}\n"
            f"Provide a valid path to a PNG, JPG, or other image file."
        )
    ext = img.suffix.lstrip(".").lower()
    if ext not in ALLOWED_IMAGE_EXTENSIONS:
        raise ValueError(
            f"Unsupported image format '.{ext}'.\n"
            f"Allowed: {', '.join(sorted(ALLOWED_IMAGE_EXTENSIONS))}"
        )

    from .validation import VALID_PAGE_SIZES
    size_key = page_size.lower()
    if size_key not in VALID_PAGE_SIZES:
        raise ValueError(
            f"Unknown page size '{page_size}'.\n"
            f"Valid sizes: {', '.join(sorted(VALID_PAGE_SIZES))}"
        )
    w_mm, h_mm = VALID_PAGE_SIZES[size_key]
    if landscape:
        w_mm, h_mm = h_mm, w_mm

    cx = int(Mm(width_mm if width_mm is not None else w_mm))
    cy = int(Mm(height_mm if height_mm is not None else h_mm))
    offset_x = int(Mm(offset_x_mm))
    offset_y = int(Mm(offset_y_mm))

    if not 1 <= opacity <= 100:
        raise ValueError(
            f"opacity={opacity} is out of range [1, 100].\n"
            f"Use 100 for fully opaque, 30-50 for a faded watermark effect."
        )

    safety.check_file_size(img)

    doc = Document(str(p))

    max_id = max(
        (int(el.get("id", 0))
         for el in doc.element.iter()
         if el.tag.endswith("}docPr") and el.get("id", "").isdigit()),
        default=0,
    )
    doc_pr_id = max_id + 1

    for section in doc.sections:
        header = section.header
        header.is_linked_to_previous = False

        _remove_bg_images_from_header(header._element)

        r_id = header.part.relate_to(
            doc.part.package.get_or_add_image_part(str(img)),
            RT.IMAGE,
        )

        xml_str = _build_bg_anchor_xml(
            r_id=r_id, cx=cx, cy=cy,
            offset_x=offset_x, offset_y=offset_y,
            doc_pr_id=doc_pr_id, opacity_pct=opacity,
        )
        anchor_el = safe_parse_xml(xml_str.encode("utf-8"))

        bg_para = etree.SubElement(header._element, f"{{{W_NS}}}p")
        bg_para.append(anchor_el)
        header._element.remove(bg_para)
        header._element.insert(0, bg_para)

        doc_pr_id += 1

    log.info(
        "Background image set: %s → %s (offset=%d,%dmm size=%dx%dmm opacity=%d%%)",
        img.name, p.name, offset_x_mm, offset_y_mm,
        width_mm or w_mm, height_mm or h_mm, opacity,
    )
    doc.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# DOCX watermark (VML in header)
# --------------------------------------------------------------------------

_WATERMARK_VML_TEMPLATE = """
<w:p xmlns:w="{w_ns}" xmlns:v="{v_ns}" xmlns:o="{o_ns}">
  <w:r>
    <w:pict>
      <v:shape id="WatermarkShape" o:spid="_x0000_s1026" type="#_x0000_t136"
        style="position:absolute;left:0;text-align:left;margin-left:0;margin-top:0;
               width:468pt;height:54pt;rotation:315;z-index:-251658240;
               mso-position-horizontal:center;mso-position-horizontal-relative:margin;
               mso-position-vertical:center;mso-position-vertical-relative:margin"
        fillcolor="#{color}" stroked="f">
        <v:textpath style="font-family:&quot;Calibri&quot;;font-size:{size}pt"
          string="{text}"/>
      </v:shape>
    </w:pict>
  </w:r>
</w:p>
"""


def _remove_watermarks_from_header(header_element: etree._Element) -> None:
    """Remove any existing WatermarkShape VML elements from a header."""
    for p_el in list(header_element.findall(f"{{{W_NS}}}p")):
        for r_el in p_el.findall(f"{{{W_NS}}}r"):
            for pict in r_el.findall(f"{{{W_NS}}}pict"):
                for shape in pict.findall(f"{{{V_NS}}}shape"):
                    if shape.get("id") == "WatermarkShape":
                        header_element.remove(p_el)
                        break


def docx_set_watermark(
    path: str,
    text: str,
    *,
    color_hex: str = "#BFBFBF",
    font_size: int = 144,
) -> str:
    """Add a diagonal text watermark to every section's header.

    Replaces any previously set watermark. Best-effort across renderers:
    Word renders fully; LibreOffice partially; Google Docs ignores VML.
    """
    from docx import Document

    color = _normalise_color(color_hex)
    p = validate_path(path, must_exist=True, expected_ext="docx", operation="set_watermark")

    if not isinstance(font_size, int):
        raise ValueError(
            f"font_size must be an integer, got {type(font_size).__name__}."
        )
    if font_size < 1 or font_size > 999:
        raise ValueError(
            f"font_size={font_size} is out of range [1, 999].\n"
            f"Common watermark sizes: 72, 96, 144."
        )
    if len(text) > 200:
        raise ValueError(
            f"Watermark text too long ({len(text)} chars, max 200)."
        )

    doc = Document(str(p))
    for section in doc.sections:
        header = section.header
        _remove_watermarks_from_header(header._element)

        safe_text = xml.sax.saxutils.escape(text, {'"': "&quot;"})
        wm_xml = _WATERMARK_VML_TEMPLATE.format(
            w_ns=W_NS, v_ns=V_NS, o_ns=O_NS,
            color=color, size=font_size, text=safe_text,
        ).strip()
        wm_bytes = wm_xml.encode("utf-8") if isinstance(wm_xml, str) else wm_xml
        element = safe_parse_xml(wm_bytes)
        header._element.append(element)

    doc.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# PPTX slide background
# --------------------------------------------------------------------------

def pptx_set_slide_background(
    path: str,
    slide_index: int | None = None,
    color_hex: str = "#FFFFFF",
) -> str:
    """Set a solid background colour on one slide (``slide_index``) or all
    slides (``slide_index=None``)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    from .validation import validate_slide_index

    color = _normalise_color(color_hex)
    rgb = RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))

    p = validate_path(path, must_exist=True, expected_ext="pptx", operation="set_slide_background")

    prs = Presentation(str(p))
    if slide_index is not None:
        validate_slide_index(prs, slide_index)
    slides = prs.slides if slide_index is None else [prs.slides[slide_index]]
    for slide in slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb
    prs.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# XLSX sheet tab colour
# --------------------------------------------------------------------------

def xlsx_set_sheet_tab_color(path: str, sheet: str, color_hex: str) -> str:
    """Set the colour of a sheet's tab strip."""
    from openpyxl import load_workbook

    from .validation import validate_sheet_name

    color = _normalise_color(color_hex)
    p = validate_path(path, must_exist=True, expected_ext="xlsx", operation="set_sheet_tab_color")

    wb = load_workbook(str(p))
    validate_sheet_name(wb, sheet)
    wb[sheet].sheet_properties.tabColor = color
    wb.save(str(p))
    return str(p)
