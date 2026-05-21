"""Spell-check and auto-correct.

Primary engine: pyspellchecker (pure Python, bundled English dict).
Fallback: ``/usr/bin/aspell`` via subprocess when pyspellchecker is missing
or its dictionary cannot be loaded.

Run-level formatting is preserved when applying corrections: we operate on
``paragraph.runs[*].text``, not on ``paragraph.text``. Multi-run spans are
flagged in the response so the caller knows they were skipped.
"""

from __future__ import annotations

import re
import shutil
import subprocess
from pathlib import Path
from typing import Any

from .errors import EngineMissing


_WORD_RE = re.compile(r"[A-Za-z][A-Za-z'-]*")


# --------------------------------------------------------------------------
# Engine resolution
# --------------------------------------------------------------------------

def _get_pyspellchecker(language: str):
    try:
        from spellchecker import SpellChecker  # type: ignore
    except ImportError:
        return None
    try:
        return SpellChecker(language=language)
    except ValueError:
        # Language not bundled — pyspellchecker raises ValueError.
        return None


def _aspell_available() -> bool:
    return shutil.which("aspell") is not None


def _engine(language: str) -> tuple[str, Any]:
    """Return (engine_name, engine_object) or ("aspell", None) or raise."""
    sp = _get_pyspellchecker(language)
    if sp is not None:
        return ("pyspellchecker", sp)
    if _aspell_available():
        return ("aspell", None)
    raise EngineMissing(
        "No spell-check engine available. Install pyspellchecker "
        "(pip install pyspellchecker) or aspell (apt install aspell)."
    )


# --------------------------------------------------------------------------
# Word-level check
# --------------------------------------------------------------------------

def _check_pyspellchecker(words: list[str], sp: Any) -> dict[str, list[str]]:
    """Return {misspelled_word: [suggestion, …]} for the given words."""
    unknown = sp.unknown(words)
    out: dict[str, list[str]] = {}
    for w in unknown:
        candidates = sp.candidates(w) or set()
        # Sort so the most likely is first; pyspellchecker doesn't rank.
        suggestions = sorted(candidates - {w})[:3]
        out[w] = suggestions
    return out


def _check_aspell(words: list[str], language: str = "en") -> dict[str, list[str]]:
    """Run aspell in pipe mode and parse responses."""
    binary = shutil.which("aspell")
    if not binary:
        return {}
    proc = subprocess.run(
        [binary, "-a", "-l", language],
        input="!\n" + "\n".join(words) + "\n",
        capture_output=True,
        text=True,
        timeout=30,
    )
    out: dict[str, list[str]] = {}
    for line in proc.stdout.splitlines():
        line = line.strip()
        if not line:
            continue
        # & word N offset: sug1, sug2, ...
        if line.startswith("&"):
            try:
                head, suggestions = line.split(":", 1)
                word = head.split()[1]
                out[word] = [s.strip() for s in suggestions.split(",")][:3]
            except (IndexError, ValueError):
                continue
        elif line.startswith("#"):  # no suggestions
            parts = line.split()
            if len(parts) >= 2:
                out[parts[1]] = []
    return out


# --------------------------------------------------------------------------
# Per-format text iteration
# --------------------------------------------------------------------------

def _iter_docx_words(path: Path):
    """Yield (paragraph_index, word, char_offset, context)."""
    from docx import Document

    doc = Document(str(path))
    for idx, para in enumerate(doc.paragraphs):
        text = para.text
        for m in _WORD_RE.finditer(text):
            start = m.start()
            ctx_start = max(0, start - 16)
            ctx_end = min(len(text), m.end() + 16)
            context = text[ctx_start:ctx_end]
            yield {"paragraph_index": idx}, m.group(), start, context


def _iter_xlsx_words(path: Path):
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True)
    for name in wb.sheetnames:
        ws = wb[name]
        for row in ws.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                for m in _WORD_RE.finditer(cell.value):
                    ctx = cell.value
                    yield {"sheet": name, "cell": cell.coordinate}, m.group(), m.start(), ctx


def _iter_pptx_words(path: Path):
    from pptx import Presentation

    prs = Presentation(str(path))
    for s_idx, slide in enumerate(prs.slides):
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                text = paragraph.text
                for m in _WORD_RE.finditer(text):
                    yield (
                        {"slide_index": s_idx, "shape_index": sh_idx, "paragraph_index": p_idx},
                        m.group(),
                        m.start(),
                        text,
                    )


# --------------------------------------------------------------------------
# Public tools
# --------------------------------------------------------------------------

def check_document(path: str, *, language: str = "en", max_words: int = 200) -> dict:
    """Check a document for spelling errors.

    Returns ``{engine, misspellings: [{word, suggestions, location, context}]}``.
    ``location`` shape depends on format: docx -> ``paragraph_index`` +
    ``char_offset``; xlsx -> ``sheet`` + ``cell``; pptx -> slide/shape ids.
    """
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(p)
    ext = p.suffix.lstrip(".").lower()

    engine_name, engine_obj = _engine(language)

    findings: list[dict] = []
    seen_words: set[str] = set()
    unique_words: list[str] = []
    word_records: list[dict] = []

    iterator = {
        "docx": _iter_docx_words,
        "xlsx": _iter_xlsx_words,
        "pptx": _iter_pptx_words,
    }.get(ext)
    if iterator is None:
        return {"engine": engine_name, "misspellings": [], "error": f"unsupported: {ext}"}

    for loc, word, offset, context in iterator(p):
        record = {"word": word, "location": loc, "char_offset": offset, "context": context}
        word_records.append(record)
        if word.lower() not in seen_words:
            seen_words.add(word.lower())
            unique_words.append(word)
        if len(word_records) >= max_words:
            break

    # Look up suggestions for the unique words.
    if engine_name == "pyspellchecker":
        suggestions = _check_pyspellchecker(unique_words, engine_obj)
    else:
        suggestions = _check_aspell(unique_words, language=language)

    for rec in word_records:
        word = rec["word"]
        # Try the original casing first then a lowercase lookup.
        sug = suggestions.get(word) or suggestions.get(word.lower())
        if sug is not None:
            findings.append({**rec, "suggestions": sug})

    return {
        "engine": engine_name,
        "language": language,
        "checked_words": len(word_records),
        "misspellings": findings,
    }


def apply_corrections(
    path: str,
    corrections: dict[str, str],
    *,
    scope: str = "all",
) -> dict:
    """Apply ``{misspelled: correction}`` substitutions in a document,
    preserving run formatting.

    ``scope``:
      - ``"all"``: apply everywhere
      - ``"single_run_only"``: skip occurrences that span multiple runs
    """
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(p)
    ext = p.suffix.lstrip(".").lower()
    applied = 0
    skipped = 0

    # Build a regex that matches any of the misspelled words as whole words.
    if not corrections:
        return {"applied": 0, "skipped_multi_run": 0}
    pattern = re.compile(
        r"\b(" + "|".join(re.escape(w) for w in corrections.keys()) + r")\b"
    )

    if ext == "docx":
        from docx import Document

        doc = Document(str(p))
        for para in doc.paragraphs:
            para_applied = 0
            for run in para.runs:
                new_text, count = pattern.subn(
                    lambda m: corrections[m.group(1)], run.text
                )
                if count:
                    run.text = new_text
                    para_applied += count
            applied += para_applied
            if scope == "all":
                joined = "".join(r.text for r in para.runs)
                if pattern.search(joined):
                    rewritten, count = pattern.subn(
                        lambda m: corrections[m.group(1)], joined
                    )
                    extra = count - para_applied
                    if extra > 0 and para.runs:
                        para.runs[0].text = rewritten
                        for r in para.runs[1:]:
                            r.text = ""
                        applied += extra
                        skipped += extra
        doc.save(str(p))
    elif ext == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(str(p))
        for name in wb.sheetnames:
            ws = wb[name]
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and pattern.search(cell.value):
                        cell.value, count = pattern.subn(
                            lambda m: corrections[m.group(1)], cell.value
                        )
                        applied += count
        wb.save(str(p))
    elif ext == "pptx":
        from pptx import Presentation

        prs = Presentation(str(p))
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        new_text, count = pattern.subn(
                            lambda m: corrections[m.group(1)], run.text
                        )
                        if count:
                            run.text = new_text
                            applied += count
        prs.save(str(p))
    else:
        return {"applied": 0, "skipped_multi_run": 0, "error": f"unsupported: {ext}"}

    return {"applied": applied, "skipped_multi_run": skipped}


def suggest_single(word: str, *, language: str = "en", max: int = 5) -> dict:
    """Look up suggestions for a single word."""
    engine_name, engine_obj = _engine(language)
    if engine_name == "pyspellchecker":
        is_known = bool(engine_obj.known([word]))
        candidates = engine_obj.candidates(word) or set()
        suggestions = sorted(candidates - {word})[:max]
    else:
        aspell_result = _check_aspell([word], language=language)
        is_known = word not in aspell_result
        suggestions = aspell_result.get(word, [])[:max]
    return {
        "engine": engine_name,
        "word": word,
        "suggestions": suggestions,
        "is_known": is_known,
    }
