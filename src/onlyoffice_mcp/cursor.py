"""AI virtual cursor — a persistent current-position hint across MCP calls."""

from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Any

import logging

from . import history
from .errors import CursorOutOfBounds

log = logging.getLogger(__name__)


_PAGE_BREAK_RE = re.compile(rb'<w:br[^>]*w:type="page"', re.IGNORECASE)
_LAST_RENDERED_RE = re.compile(rb'<w:lastRenderedPageBreak[^>]*/?>', re.IGNORECASE)


def estimate_docx_pages(path: str | Path) -> int:
    """Approximate page count by counting page breaks in document.xml."""
    p = Path(path).expanduser()
    if not p.exists():
        return 0
    try:
        with zipfile.ZipFile(str(p)) as zf:
            try:
                doc_xml = zf.read("word/document.xml")
            except KeyError:
                return 1
    except zipfile.BadZipFile:
        return 1
    explicit = len(_PAGE_BREAK_RE.findall(doc_xml))
    rendered = len(_LAST_RENDERED_RE.findall(doc_xml))
    return max(1, 1 + max(explicit, rendered))


def precise_docx_pages(path: str | Path, *, timeout: int = 60) -> int:
    from . import libreoffice
    return libreoffice.page_count_via_pdf(Path(path), timeout=timeout)


_CELL_RE = re.compile(r"^([A-Z]+)(\d+)$")


def cell_to_rc(cell: str) -> tuple[int, int]:
    m = _CELL_RE.match(cell.upper())
    if not m:
        raise ValueError(f"Invalid cell reference: {cell!r}")
    col_letters, row_str = m.groups()
    col = 0
    for ch in col_letters:
        col = col * 26 + (ord(ch) - ord("A") + 1)
    return int(row_str), col


def rc_to_cell(row: int, col: int) -> str:
    if row < 1 or col < 1:
        raise ValueError(f"Row/col must be >= 1, got row={row}, col={col}")
    letters = ""
    while col > 0:
        col, rem = divmod(col - 1, 26)
        letters = chr(ord("A") + rem) + letters
    return f"{letters}{row}"


def _doc_format(path: str | Path) -> str:
    return Path(path).suffix.lstrip(".").lower()


def _doc_bounds(path: Path) -> dict[str, Any]:
    ext = path.suffix.lstrip(".").lower()
    if ext == "docx":
        try:
            from docx import Document
            doc = Document(str(path))
            return {
                "paragraph_count": len(doc.paragraphs),
                "page_estimate": estimate_docx_pages(path),
            }
        except Exception as exc:
            log.warning("docx bounds failed for %s: %s", path, exc)
            return {"paragraph_count": 0, "page_estimate": 1}
    if ext == "xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True)
            sheets = {}
            for name in wb.sheetnames:
                ws = wb[name]
                sheets[name] = {"rows": ws.max_row or 1, "cols": ws.max_column or 1}
            return {"sheets": sheets}
        except Exception as exc:
            log.warning("xlsx bounds failed for %s: %s", path, exc)
            return {"sheets": {}}
    if ext == "pptx":
        try:
            from pptx import Presentation
            return {"slide_count": len(Presentation(str(path)).slides)}
        except Exception as exc:
            log.warning("pptx bounds failed for %s: %s", path, exc)
            return {"slide_count": 0}
    return {}


def get_cursor(path: str | Path) -> dict:
    p = Path(path).expanduser().resolve()
    meta = history._load_meta(p)
    cursor = dict(meta.get("cursor") or {})
    bounds = _doc_bounds(p)
    stale = False
    ext = _doc_format(p)

    if ext == "docx":
        para_count = bounds.get("paragraph_count", 0)
        page_est = bounds.get("page_estimate", 1)
        if cursor.get("paragraph_index", 0) >= max(1, para_count):
            stale = True
        cursor.setdefault("paragraph_index", 0)
        cursor.setdefault("page_estimate", page_est)
    elif ext == "xlsx":
        sheet = cursor.get("sheet")
        if sheet and sheet not in bounds.get("sheets", {}):
            stale = True
        if not sheet and bounds.get("sheets"):
            sheet = next(iter(bounds["sheets"]))
        cursor.setdefault("sheet", sheet or "")
        cursor.setdefault("row", 1)
        cursor.setdefault("col", 1)
    elif ext == "pptx":
        slide_count = bounds.get("slide_count", 0)
        if cursor.get("slide_index", 0) >= max(1, slide_count):
            stale = True
        cursor.setdefault("slide_index", 0)

    return {"format": ext, "cursor": cursor, "bounds": bounds, "stale": stale}


def set_cursor(
    path: str | Path,
    *,
    paragraph_index: int | None = None,
    page: int | None = None,
    sheet: str | None = None,
    cell: str | None = None,
    row: int | None = None,
    col: int | None = None,
    slide_index: int | None = None,
    clamp: bool = True,
) -> dict:
    p = Path(path).expanduser().resolve()
    meta = history._load_meta(p)
    cursor = dict(meta.get("cursor") or {})
    bounds = _doc_bounds(p)
    warnings: list[str] = []
    ext = _doc_format(p)

    def _clamp(name: str, value: int, lo: int, hi: int) -> int:
        if hi < lo:
            hi = lo
        if value < lo or value > hi:
            if not clamp:
                raise CursorOutOfBounds(f"{name}={value} out of bounds [{lo}, {hi}]")
            new = max(lo, min(value, hi))
            warnings.append(f"{name} clamped from {value} to {new}")
            return new
        return value

    if ext == "docx":
        para_max = max(0, bounds.get("paragraph_count", 1) - 1)
        page_max = max(1, bounds.get("page_estimate", 1))
        if paragraph_index is not None:
            cursor["paragraph_index"] = _clamp(
                "paragraph_index", int(paragraph_index), 0, para_max
            )
        if page is not None:
            cursor["page_estimate"] = _clamp("page", int(page), 1, page_max)
        cursor.setdefault("paragraph_index", 0)
        cursor.setdefault("page_estimate", page_max)
    elif ext == "xlsx":
        sheets = bounds.get("sheets", {})
        if sheet is not None:
            if sheets and sheet not in sheets:
                if clamp:
                    warnings.append(
                        f"sheet '{sheet}' not found; using {next(iter(sheets))}"
                    )
                    sheet = next(iter(sheets))
                else:
                    raise CursorOutOfBounds(f"sheet not found: {sheet}")
            cursor["sheet"] = sheet
        target_sheet = cursor.get("sheet") or (next(iter(sheets)) if sheets else "")
        cursor["sheet"] = target_sheet
        if cell is not None:
            row, col = cell_to_rc(cell)
        s_bounds = sheets.get(target_sheet, {"rows": 1, "cols": 1})
        if row is not None:
            cursor["row"] = _clamp("row", int(row), 1, s_bounds["rows"])
        if col is not None:
            cursor["col"] = _clamp("col", int(col), 1, s_bounds["cols"])
        cursor.setdefault("row", 1)
        cursor.setdefault("col", 1)
    elif ext == "pptx":
        slide_max = max(0, bounds.get("slide_count", 1) - 1)
        if slide_index is not None:
            cursor["slide_index"] = _clamp(
                "slide_index", int(slide_index), 0, slide_max
            )
        cursor.setdefault("slide_index", 0)

    meta["cursor"] = cursor
    history._save_meta(p, meta)
    return {"format": ext, "cursor": cursor, "warnings": warnings}


_AUTO_ADVANCE_OPS = {
    "docx_append",
    "docx_create",
    "xlsx_append_rows",
    "xlsx_create",
    "pptx_add_slide",
    "pptx_create",
}


def maybe_auto_advance(path: str | Path, tool: str) -> None:
    if tool not in _AUTO_ADVANCE_OPS:
        return
    p = Path(path).expanduser().resolve()
    ext = _doc_format(p)
    try:
        bounds = _doc_bounds(p)
        meta = history._load_meta(p)
        cursor = dict(meta.get("cursor") or {})
        if ext == "docx":
            cursor["paragraph_index"] = max(0, bounds.get("paragraph_count", 1) - 1)
            cursor["page_estimate"] = bounds.get("page_estimate", 1)
        elif ext == "xlsx":
            sheets = bounds.get("sheets", {})
            target = cursor.get("sheet") or (next(iter(sheets)) if sheets else "")
            if target in sheets:
                cursor["sheet"] = target
                cursor["row"] = sheets[target]["rows"]
                cursor["col"] = 1
        elif ext == "pptx":
            cursor["slide_index"] = max(0, bounds.get("slide_count", 1) - 1)
        meta["cursor"] = cursor
        history._save_meta(p, meta)
    except Exception as exc:
        log.warning("auto-advance failed for %s (%s): %s", path, tool, exc)


def page_count(path: str | Path, precise: bool = False) -> dict:
    p = Path(path).expanduser().resolve()
    ext = _doc_format(p)
    if ext == "docx":
        approx = estimate_docx_pages(p)
        result: dict[str, Any] = {
            "format": "docx",
            "page_count_estimate": approx,
            "method": "page_breaks + lastRenderedPageBreak",
            "note": "DOCX page count is layout-dependent; this is approximate.",
        }
        if precise:
            try:
                result["page_count_precise"] = precise_docx_pages(p)
                result["method"] = "libreoffice -> pdf -> pypdf"
            except Exception as e:
                result["precise_error"] = str(e)
        return result
    if ext == "xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(p), read_only=True)
            return {
                "format": "xlsx",
                "sheet_count": len(wb.sheetnames),
                "sheet_names": list(wb.sheetnames),
            }
        except Exception as e:
            return {"format": "xlsx", "error": str(e)}
    if ext == "pptx":
        try:
            from pptx import Presentation
            return {
                "format": "pptx",
                "slide_count": len(Presentation(str(p)).slides),
            }
        except Exception as e:
            return {"format": "pptx", "error": str(e)}
    return {"format": ext, "error": "unsupported format"}
