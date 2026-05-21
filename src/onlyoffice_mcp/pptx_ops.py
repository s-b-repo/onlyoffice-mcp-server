"""PowerPoint operations using python-pptx."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from .validation import validate_path, validate_slide_def, sanitize_text, validate_color


# python-pptx ships with the default Office template. Layout indexes:
#   0 = Title Slide        (title + subtitle)
#   1 = Title and Content  (title + body bullets)
#   5 = Title Only         (title at top, rest blank)
#   6 = Blank
_LAYOUT_INDEX = {
    "title": 0,
    "content": 1,
    "title_only": 5,
    "blank": 6,
    "image": 5,
}


def _layout(prs: Presentation, name: str):
    idx = _LAYOUT_INDEX.get(name, 1)
    # Be defensive — some templates have fewer layouts than expected.
    if idx >= len(prs.slide_layouts):
        idx = min(idx, len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


def _set_body_bullets(placeholder, bullets: list[str]) -> None:
    if not bullets:
        return
    tf = placeholder.text_frame
    tf.text = sanitize_text(str(bullets[0]), "slide bullet")
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = sanitize_text(str(line), "slide bullet")


def create(path: str, slides: list[dict]) -> str:
    """Create a .pptx file at `path`.

    Each slide dict supports:
      - {"layout": "title", "title": str, "subtitle": str}
      - {"layout": "content", "title": str, "body": [bullets]}
      - {"layout": "title_only", "title": str}
      - {"layout": "image", "title": str, "image_path": str,
         "left_inches": float, "top_inches": float, "width_inches": float}
      - {"layout": "blank"}
    """
    out = validate_path(path, expected_ext="pptx", for_creation=True, operation="create")
    prs = Presentation()

    for slide_def in slides:
        validate_slide_def(slide_def)
        layout_name = slide_def.get("layout", "content")
        slide = prs.slides.add_slide(_layout(prs, layout_name))

        # Title is always placeholder index 0 (when present).
        if slide.shapes.title and "title" in slide_def:
            slide.shapes.title.text = sanitize_text(str(slide_def["title"]), "slide title")

        if layout_name == "title":
            if len(slide.placeholders) > 1 and "subtitle" in slide_def:
                slide.placeholders[1].text = sanitize_text(str(slide_def["subtitle"]), "slide subtitle")

        elif layout_name == "content":
            body_items = slide_def.get("body", [])
            if body_items and len(slide.placeholders) > 1:
                _set_body_bullets(slide.placeholders[1], body_items)

        elif layout_name == "image":
            img = slide_def.get("image_path")
            if img:
                img_p = Path(img).expanduser().resolve()
                if not img_p.exists():
                    raise ValueError(f"Image not found: {img_p}")
                left = Inches(slide_def.get("left_inches", 1.0))
                top = Inches(slide_def.get("top_inches", 1.5))
                width = Inches(slide_def.get("width_inches", 8.0))
                slide.shapes.add_picture(str(img_p), left, top, width=width)

        # Add a notes block if requested.
        if slide_def.get("notes"):
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = sanitize_text(str(slide_def["notes"]), "speaker notes")

    prs.save(str(out))
    return str(out)


def read(path: str) -> dict:
    """Extract slide titles and body text from a .pptx file."""
    in_ = validate_path(path, must_exist=True, expected_ext="pptx", operation="read")
    prs = Presentation(str(in_))
    result: dict[str, Any] = {"slide_count": len(prs.slides), "slides": []}
    for i, slide in enumerate(prs.slides):
        info: dict[str, Any] = {"index": i, "title": "", "text": [], "notes": ""}
        title_shape = slide.shapes.title
        title_el = title_shape._element if title_shape is not None else None
        if title_shape and title_shape.has_text_frame:
            info["title"] = title_shape.text_frame.text
        for shape in slide.shapes:
            if title_el is not None and shape._element is title_el:
                continue
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text:
                    info["text"].append(paragraph.text)
        if slide.has_notes_slide:
            info["notes"] = slide.notes_slide.notes_text_frame.text
        result["slides"].append(info)
    return result


def add_slide(path: str, slide_def: dict) -> str:
    """Append a single slide to an existing .pptx file."""
    validate_slide_def(slide_def)
    in_ = validate_path(path, must_exist=True, expected_ext="pptx", operation="add_slide")
    prs = Presentation(str(in_))

    layout_name = slide_def.get("layout", "content")
    slide = prs.slides.add_slide(_layout(prs, layout_name))

    if slide.shapes.title and "title" in slide_def:
        slide.shapes.title.text = sanitize_text(str(slide_def["title"]), "slide title")
    if layout_name == "title" and len(slide.placeholders) > 1 and "subtitle" in slide_def:
        slide.placeholders[1].text = sanitize_text(str(slide_def["subtitle"]), "slide subtitle")
    elif layout_name == "content":
        body_items = slide_def.get("body", [])
        if body_items and len(slide.placeholders) > 1:
            _set_body_bullets(slide.placeholders[1], body_items)
    elif layout_name == "image" and slide_def.get("image_path"):
        img_p = Path(slide_def["image_path"]).expanduser().resolve()
        if not img_p.exists():
            raise ValueError(f"Image not found: {img_p}")
        slide.shapes.add_picture(
            str(img_p),
            Inches(slide_def.get("left_inches", 1.0)),
            Inches(slide_def.get("top_inches", 1.5)),
            width=Inches(slide_def.get("width_inches", 8.0)),
        )

    if slide_def.get("notes"):
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = sanitize_text(str(slide_def["notes"]), "speaker notes")

    prs.save(str(in_))
    return str(in_)


def update_slide(
    path: str,
    slide_index: int,
    *,
    title: str | None = None,
    body: list[str] | None = None,
    notes: str | None = None,
) -> dict:
    """Edit an existing slide's title, body text, or speaker notes in place.

    Only the parameters you provide are changed — everything else is preserved.
    Returns the slide's state after editing.
    """
    from .validation import validate_slide_index

    in_ = validate_path(path, must_exist=True, expected_ext="pptx", operation="update_slide")
    prs = Presentation(str(in_))
    validate_slide_index(prs, slide_index)
    slide = prs.slides[slide_index]

    if title is not None:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide.shapes.title.text = sanitize_text(str(title), "slide title")
        else:
            raise ValueError(
                f"Slide {slide_index} has no title placeholder.\n"
                f"Use a slide layout with a title (e.g. 'content' or 'title')."
            )

    if body is not None:
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body_ph = ph
                break
        if body_ph is None or not body_ph.has_text_frame:
            raise ValueError(
                f"Slide {slide_index} has no body placeholder.\n"
                f"Body text can only be set on slides with a 'content' layout."
            )
        _set_body_bullets(body_ph, body)

    if notes is not None:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = sanitize_text(str(notes), "speaker notes")

    prs.save(str(in_))

    result: dict = {
        "path": str(in_),
        "slide_index": slide_index,
        "title": "",
        "body": [],
        "notes": "",
    }
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        result["title"] = slide.shapes.title.text
    title_el = slide.shapes.title._element if slide.shapes.title else None
    for shape in slide.shapes:
        if title_el is not None and shape._element is title_el:
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text:
                    result["body"].append(para.text)
    if slide.has_notes_slide:
        result["notes"] = slide.notes_slide.notes_text_frame.text
    return result


def delete_slide(path: str, slide_index: int) -> str:
    """Delete a slide by 0-based index from a .pptx file."""
    from .validation import validate_slide_index

    in_ = validate_path(path, must_exist=True, expected_ext="pptx", operation="delete_slide")
    prs = Presentation(str(in_))
    validate_slide_index(prs, slide_index)
    if len(prs.slides) <= 1:
        raise ValueError(
            "Cannot delete the last slide in a presentation.\n"
            "A presentation must have at least one slide."
        )
    rId = prs.slides._sldIdLst[slide_index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    if rId is None:
        raise ValueError(
            f"Could not find relationship ID for slide {slide_index}.\n"
            f"The presentation's internal structure may be corrupted."
        )
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]
    prs.save(str(in_))
    return str(in_)
