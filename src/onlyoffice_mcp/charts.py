"""Chart insertion across docx / xlsx / pptx.

- **docx** -> render with matplotlib (Agg backend), save to a temp PNG,
  embed via ``python-docx.add_picture``. Result is a static image — NOT an
  editable chart. Documented in the tool docstring.
- **xlsx** -> native ``openpyxl`` charts. Editable inside the workbook.
- **pptx** -> native ``python-pptx`` charts. Editable inside the deck.

Force ``matplotlib.use("Agg")`` BEFORE any pyplot import so that chart
rendering works in headless / stdio-only contexts.
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Any

# Matplotlib must use a non-GUI backend.
import matplotlib

matplotlib.use("Agg")  # noqa: E402
import matplotlib.pyplot as plt  # noqa: E402


# --------------------------------------------------------------------------
# Supported chart kinds + synonyms
# --------------------------------------------------------------------------

_KINDS = ("bar", "line", "pie", "scatter", "area")
_SYNONYMS = {
    "column": "bar",
    "histogram": "bar",
    "doughnut": "pie",
    "donut": "pie",
    "xy": "scatter",
    "scatterplot": "scatter",
    "stacked_bar": "bar",
}


def _normalise_kind(kind: str) -> str:
    from .validation import validate_chart_type
    return validate_chart_type(kind)


def chart_kinds_info() -> dict:
    """Return supported chart kinds + synonyms + per-format support."""
    return {
        "kinds": list(_KINDS),
        "synonyms": dict(_SYNONYMS),
        "supported_in": {
            "docx": list(_KINDS),
            "xlsx": list(_KINDS),
            "pptx": ["bar", "line", "pie", "scatter"],  # area available but XL_CHART_TYPE varies
        },
    }


# --------------------------------------------------------------------------
# DOCX (matplotlib -> PNG -> add_picture)
# --------------------------------------------------------------------------

def docx_add_chart(
    path: str,
    chart_type: str,
    categories: list[Any],
    series: list[dict],
    *,
    title: str | None = None,
    xlabel: str | None = None,
    ylabel: str | None = None,
    width_inches: float = 6.0,
    height_inches: float = 4.0,
    paragraph_index: int | None = None,
    data_labels: bool = False,
    legend_position: str = "best",
    colors: list[str] | None = None,
    stacked: bool = False,
    horizontal: bool = False,
    explode: list[float] | None = None,
    donut: float | None = None,
    line_styles: list[str] | None = None,
    grid: bool = True,
    dpi: int = 150,
) -> str:
    """Render a chart with matplotlib and embed it as a picture in a .docx.

    ``series`` is a list of ``{"name": str, "values": [numbers]}`` for
    bar/line/area/pie, or ``{"name": str, "x": [numbers], "y": [numbers]}``
    for scatter.

    Note: the embedded result is a static image — the chart cannot be edited
    inside Word/OnlyOffice. Use ``xlsx_add_chart`` / ``pptx_add_chart`` for
    editable charts.
    """
    from docx import Document
    from docx.shared import Inches

    from .validation import validate_color, validate_series_data

    kind = _normalise_kind(chart_type)
    validate_series_data(series, kind)
    out = Path(path).expanduser().resolve()
    if not out.exists():
        raise FileNotFoundError(out)

    color_list = [f"#{validate_color(c)}" for c in colors] if colors else None
    leg_loc = legend_position or "best"

    fig, ax = plt.subplots(figsize=(width_inches, height_inches), dpi=dpi)

    if kind == "pie":
        values = list(series[0].get("values", []))
        labels = [str(c) for c in categories[: len(values)]]
        pie_colors = color_list[:len(values)] if color_list else None
        explode_vals = explode[:len(values)] if explode else None
        wedgeprops = dict(width=1 - max(0.0, min(donut, 0.8))) if donut else None
        autopct = "%1.1f%%" if data_labels else None
        ax.pie(values, labels=labels, autopct=autopct, startangle=90,
               colors=pie_colors, explode=explode_vals,
               **({"wedgeprops": wedgeprops} if wedgeprops else {}))
        ax.axis("equal")

    elif kind == "scatter":
        for i, s in enumerate(series):
            c = color_list[i % len(color_list)] if color_list else None
            ax.scatter(s.get("x", []), s.get("y", []), label=s.get("name"), color=c)
            if data_labels:
                for j, (xv, yv) in enumerate(zip(s.get("x", []), s.get("y", []))):
                    ax.annotate(f"({xv},{yv})", (xv, yv), textcoords="offset points",
                                xytext=(5, 5), fontsize=7)
        ax.legend(loc=leg_loc)

    elif kind == "line":
        for i, s in enumerate(series):
            c = color_list[i % len(color_list)] if color_list else None
            ls = line_styles[i % len(line_styles)] if line_styles else "-"
            values = s.get("values", [])
            ax.plot(categories, values, label=s.get("name"), marker="o", color=c, linestyle=ls)
            if data_labels:
                for j, v in enumerate(values):
                    ax.annotate(f"{v}", (categories[j], v), textcoords="offset points",
                                xytext=(0, 8), ha="center", fontsize=8)
        ax.legend(loc=leg_loc)

    elif kind == "area":
        bottom = [0] * len(categories)
        for i, s in enumerate(series):
            values = list(s.get("values", []))
            c = color_list[i % len(color_list)] if color_list else None
            top = [b + v for b, v in zip(bottom, values)]
            ax.fill_between(categories, bottom, top, label=s.get("name"), alpha=0.6, color=c)
            if data_labels:
                for j, v in enumerate(values):
                    mid = bottom[j] + v / 2
                    ax.annotate(f"{v}", (categories[j], mid), ha="center", fontsize=8)
            bottom = top
        ax.legend(loc=leg_loc)

    else:  # bar
        n = len(series)
        x_positions = list(range(len(categories)))
        if stacked:
            bottom = [0.0] * len(categories)
            for i, s in enumerate(series):
                values = s.get("values", [])
                c = color_list[i % len(color_list)] if color_list else None
                if horizontal:
                    bars = ax.barh(x_positions, values, left=bottom, label=s.get("name"), color=c)
                else:
                    bars = ax.bar(x_positions, values, bottom=bottom, label=s.get("name"), color=c)
                if data_labels:
                    ax.bar_label(bars, fmt="%.0f", label_type="center", fontsize=8)
                bottom = [b + v for b, v in zip(bottom, values)]
        else:
            width = 0.8 / max(1, n)
            for i, s in enumerate(series):
                offsets = [x + (i - (n - 1) / 2) * width for x in x_positions]
                c = color_list[i % len(color_list)] if color_list else None
                if horizontal:
                    bars = ax.barh(offsets, s.get("values", []), height=width,
                                   label=s.get("name"), color=c)
                else:
                    bars = ax.bar(offsets, s.get("values", []), width=width,
                                  label=s.get("name"), color=c)
                if data_labels:
                    ax.bar_label(bars, fmt="%.0f", fontsize=8)
        if horizontal:
            ax.set_yticks(x_positions)
            ax.set_yticklabels([str(c) for c in categories])
        else:
            ax.set_xticks(x_positions)
            ax.set_xticklabels([str(c) for c in categories])
        ax.legend(loc=leg_loc)

    if title:
        ax.set_title(title)
    if kind != "pie":
        if xlabel:
            ax.set_xlabel(xlabel)
        if ylabel:
            ax.set_ylabel(ylabel)
        ax.grid(grid)
    fig.tight_layout()

    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        fig.savefig(tmp_path, bbox_inches="tight", dpi=dpi)
        plt.close(fig)

        doc = Document(str(out))
        if paragraph_index is not None and 0 <= paragraph_index < len(doc.paragraphs):
            target = doc.paragraphs[paragraph_index]
            doc.add_picture(tmp_path, width=Inches(width_inches))
            pic_para = doc.element.body[-1]
            doc.element.body.remove(pic_para)
            target._p.addnext(pic_para)
        else:
            doc.add_picture(tmp_path, width=Inches(width_inches))
        doc.save(str(out))
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    return str(out)


# --------------------------------------------------------------------------
# XLSX (native openpyxl)
# --------------------------------------------------------------------------

_XLSX_LEGEND_POS = {
    "bottom": "b", "top": "t", "left": "l", "right": "r",
    "top_right": "tr",
}


def xlsx_add_chart(
    path: str,
    sheet: str,
    chart_type: str,
    data_range: str,
    *,
    categories_range: str | None = None,
    anchor_cell: str = "E2",
    title: str | None = None,
    xlabel: str | None = None,
    ylabel: str | None = None,
    stacked: bool = False,
    legend_position: str | None = None,
    data_labels: bool = False,
) -> str:
    """Add a native chart to an existing .xlsx file.

    ``data_range`` and ``categories_range`` are A1-style refs *within the
    sheet*: e.g. ``"B1:B10"`` (first row is treated as series title). If
    ``categories_range`` is omitted, the first column of ``data_range`` is
    used as categories.
    """
    from openpyxl import load_workbook
    from openpyxl.chart import (
        AreaChart,
        BarChart,
        LineChart,
        PieChart,
        Reference,
        ScatterChart,
    )

    from .validation import validate_cell_range, validate_sheet_name

    kind = _normalise_kind(chart_type)
    validate_cell_range(data_range, "data_range")
    if categories_range:
        validate_cell_range(categories_range, "categories_range")
    out = Path(path).expanduser().resolve()
    if not out.exists():
        raise FileNotFoundError(out)

    wb = load_workbook(str(out))
    validate_sheet_name(wb, sheet)
    ws = wb[sheet]

    if ws[anchor_cell].value is None:
        ws[anchor_cell] = None

    chart_cls = {
        "bar": BarChart,
        "line": LineChart,
        "pie": PieChart,
        "scatter": ScatterChart,
        "area": AreaChart,
    }[kind]
    chart = chart_cls()
    if title:
        chart.title = title
    if stacked and kind in ("bar", "area"):
        chart.grouping = "stacked"
    if xlabel and hasattr(chart, "x_axis"):
        chart.x_axis.title = xlabel
    if ylabel and hasattr(chart, "y_axis"):
        chart.y_axis.title = ylabel
    if legend_position:
        pos = _XLSX_LEGEND_POS.get(legend_position.lower(), legend_position)
        chart.legend.position = pos
    if data_labels:
        from openpyxl.chart.label import DataLabelList
        chart.dataLabels = DataLabelList()
        chart.dataLabels.showVal = True

    chart.add_data(
        Reference(ws, range_string=f"{sheet}!{data_range}"),
        titles_from_data=True,
    )
    if categories_range:
        chart.set_categories(Reference(ws, range_string=f"{sheet}!{categories_range}"))

    ws.add_chart(chart, anchor_cell)
    wb.save(str(out))
    return str(out)


# --------------------------------------------------------------------------
# PPTX (native python-pptx)
# --------------------------------------------------------------------------

_PPTX_LEGEND_POS = {}  # populated lazily to avoid import at module level


def pptx_add_chart(
    path: str,
    slide_index: int,
    chart_type: str,
    categories: list[Any],
    series: list[dict],
    *,
    left_inches: float = 1.0,
    top_inches: float = 2.0,
    width_inches: float = 8.0,
    height_inches: float = 5.0,
    title: str | None = None,
    data_labels: bool = False,
    legend_position: str | None = None,
    stacked: bool = False,
) -> str:
    """Add a native chart to an existing slide in a .pptx file.

    For bar / line / pie / area: ``series`` is a list of ``{"name", "values"}``.
    For scatter: ``series`` is a list of ``{"name", "x", "y"}``.
    """
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    from .validation import validate_series_data

    kind = _normalise_kind(chart_type)
    validate_series_data(series, kind)
    kind_map = {
        "bar": XL_CHART_TYPE.BAR_STACKED if stacked else XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "area": XL_CHART_TYPE.AREA_STACKED if stacked else XL_CHART_TYPE.AREA,
    }
    if kind not in kind_map:
        raise ValueError(f"Chart kind '{kind}' not supported in pptx")

    out = Path(path).expanduser().resolve()
    if not out.exists():
        raise FileNotFoundError(out)

    prs = Presentation(str(out))
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"slide_index {slide_index} out of range [0, {len(prs.slides) - 1}]")
    slide = prs.slides[slide_index]

    if kind == "scatter":
        chart_data = XyChartData()
        for s in series:
            ser = chart_data.add_series(s.get("name", "Series"))
            xs = s.get("x", [])
            ys = s.get("y", [])
            for x_val, y_val in zip(xs, ys):
                ser.add_data_point(x_val, y_val)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]
        for s in series:
            chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    chart_shape = slide.shapes.add_chart(
        kind_map[kind],
        Inches(left_inches),
        Inches(top_inches),
        Inches(width_inches),
        Inches(height_inches),
        chart_data,
    )
    chart_obj = chart_shape.chart
    if title:
        chart_obj.has_title = True
        chart_obj.chart_title.text_frame.text = title
    if data_labels and chart_obj.plots:
        plot = chart_obj.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_value = True
    if legend_position:
        from pptx.enum.chart import XL_LEGEND_POSITION
        _pos_map = {
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "left": XL_LEGEND_POSITION.LEFT,
            "right": XL_LEGEND_POSITION.RIGHT,
            "top": XL_LEGEND_POSITION.TOP,
            "top_right": XL_LEGEND_POSITION.CORNER,
        }
        pos = _pos_map.get(legend_position.lower())
        if pos is not None:
            chart_obj.has_legend = True
            chart_obj.legend.position = pos

    prs.save(str(out))
    return str(out)
