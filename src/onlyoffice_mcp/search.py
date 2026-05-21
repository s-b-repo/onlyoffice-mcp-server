"""Find & replace across docx / xlsx / pptx with run-level formatting
preservation.

Implementation note: in docx and pptx, a "match" may span multiple runs
inside a paragraph. The naive ``run.text.replace`` then misses it. Strategy:
  1. Join all runs into one text string.
  2. Run the replace on the joined text.
  3. If a replacement happened that wasn't already covered by per-run
     replacements, rewrite the paragraph via the first run (carrying its
     formatting) and empty the remaining runs.

Multi-run rewrites lose fine-grained per-character formatting inside the
matched span; these are reported back as ``skipped_multi_run`` so the
caller knows.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any


def _compile_pattern(pattern: str, *, regex: bool, case_sensitive: bool) -> re.Pattern:
    if not pattern:
        raise ValueError(
            "Search pattern cannot be empty.\n"
            "Provide a string to search for, or a regex pattern with regex=True."
        )
    flags = 0 if case_sensitive else re.IGNORECASE
    if not regex:
        pattern = re.escape(pattern)
        return re.compile(pattern, flags)
    from .validation import validate_regex
    validate_regex(pattern)
    return re.compile(pattern, flags)


# --------------------------------------------------------------------------
# Find
# --------------------------------------------------------------------------

def find_in_document(
    path: str,
    pattern: str,
    *,
    regex: bool = False,
    case_sensitive: bool = True,
    include_notes: bool = True,
) -> list[dict]:
    """Find every match in a document. Per-format location:

      docx:  {"paragraph_index": int, "char_offset": int}
      xlsx:  {"sheet": str, "cell": str}
      pptx:  {"slide_index": int, "shape_index": int, "paragraph_index": int,
              "char_offset": int}
    """
    from .validation import validate_path
    p = validate_path(path, must_exist=True, operation="find")
    ext = p.suffix.lstrip(".").lower()
    pat = _compile_pattern(pattern, regex=regex, case_sensitive=case_sensitive)
    results: list[dict] = []

    def _add(loc: dict, m: re.Match, ctx: str) -> None:
        ctx_start = max(0, m.start() - 32)
        ctx_end = min(len(ctx), m.end() + 32)
        results.append(
            {
                "location": loc,
                "match_text": m.group(0),
                "context": ctx[ctx_start:ctx_end],
            }
        )

    if ext == "docx":
        from docx import Document

        doc = Document(str(p))
        for idx, para in enumerate(doc.paragraphs):
            for m in pat.finditer(para.text):
                _add({"paragraph_index": idx, "char_offset": m.start()}, m, para.text)
        for t_idx, table in enumerate(doc.tables):
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    for m in pat.finditer(cell.text):
                        _add(
                            {
                                "table_index": t_idx,
                                "row": r_idx,
                                "col": c_idx,
                                "char_offset": m.start(),
                            },
                            m,
                            cell.text,
                        )
    elif ext == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(str(p), data_only=True)
        for name in wb.sheetnames:
            ws = wb[name]
            for row in ws.iter_rows():
                for cell in row:
                    if not isinstance(cell.value, str):
                        continue
                    for m in pat.finditer(cell.value):
                        _add({"sheet": name, "cell": cell.coordinate}, m, cell.value)
    elif ext == "pptx":
        from pptx import Presentation

        prs = Presentation(str(p))
        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for p_idx, para in enumerate(shape.text_frame.paragraphs):
                    for m in pat.finditer(para.text):
                        _add(
                            {
                                "slide_index": s_idx,
                                "shape_index": sh_idx,
                                "paragraph_index": p_idx,
                                "char_offset": m.start(),
                            },
                            m,
                            para.text,
                        )
            if include_notes and slide.has_notes_slide:
                for p_idx, para in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                    for m in pat.finditer(para.text):
                        _add(
                            {
                                "slide_index": s_idx,
                                "in_notes": True,
                                "paragraph_index": p_idx,
                                "char_offset": m.start(),
                            },
                            m,
                            para.text,
                        )
    else:
        raise ValueError(
            f"Unsupported format for search: '.{ext}'.\n"
            f"Supported formats: docx, xlsx, pptx."
        )
    return results


# --------------------------------------------------------------------------
# Replace
# --------------------------------------------------------------------------

def _replace_in_paragraph_runs(para, pat: re.Pattern, replace: str) -> tuple[int, int]:
    """Replace ``pat`` -> ``replace`` in a paragraph, preserving run
    formatting where possible.

    Returns (replacements_applied, replacements_skipped_multi_run).
    """
    applied = 0
    skipped = 0

    original_joined = "".join(r.text for r in para.runs)

    # Per-run pass — covers in-run matches with formatting preservation.
    for run in para.runs:
        new_text, count = pat.subn(replace, run.text)
        if count:
            run.text = new_text
            applied += count

    # Span-crossing matches: check the ORIGINAL joined text for matches
    # that were not found in any individual run.
    if pat.search(original_joined):
        _, total_in_original = pat.subn(replace, original_joined)
        extra = max(0, total_in_original - applied)
        if extra > 0 and para.runs:
            full_replaced, _ = pat.subn(replace, original_joined)
            para.runs[0].text = full_replaced
            for r in para.runs[1:]:
                r.text = ""
            applied = total_in_original
            skipped += extra

    return applied, skipped


def replace_in_document(
    path: str,
    find: str,
    replace: str,
    *,
    regex: bool = False,
    case_sensitive: bool = True,
    count: int | None = None,
    dry_run: bool = False,
    include_formulas: bool = False,
    include_notes: bool = True,
) -> dict:
    """Replace ``find`` with ``replace`` in a document. Returns counts."""
    from .validation import validate_path
    p = validate_path(path, must_exist=True, operation="replace")
    ext = p.suffix.lstrip(".").lower()
    pat = _compile_pattern(find, regex=regex, case_sensitive=case_sensitive)

    if dry_run:
        matches = find_in_document(
            str(p),
            find,
            regex=regex,
            case_sensitive=case_sensitive,
            include_notes=include_notes,
        )
        return {
            "replacements_made": 0,
            "would_replace": len(matches),
            "skipped_multi_run": 0,
            "locations": [m["location"] for m in matches[: (count or len(matches))]],
            "dry_run": True,
        }

    total_applied = 0
    total_skipped = 0
    locations: list[dict] = []
    remaining = count if count is not None else float("inf")

    if ext == "docx":
        from docx import Document

        doc = Document(str(p))
        for idx, para in enumerate(doc.paragraphs):
            if remaining <= 0:
                break
            a, s = _replace_in_paragraph_runs(para, pat, replace)
            if a:
                total_applied += min(a, remaining)
                total_skipped += s
                locations.append({"paragraph_index": idx, "count": a})
                remaining -= a
        for t_idx, table in enumerate(doc.tables):
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    if remaining <= 0:
                        break
                    for para in cell.paragraphs:
                        a, s = _replace_in_paragraph_runs(para, pat, replace)
                        if a:
                            total_applied += min(a, remaining)
                            total_skipped += s
                            locations.append(
                                {"table_index": t_idx, "row": r_idx, "col": c_idx, "count": a}
                            )
                            remaining -= a
        doc.save(str(p))
    elif ext == "xlsx":
        from openpyxl import load_workbook

        wb = load_workbook(str(p))
        for name in wb.sheetnames:
            ws = wb[name]
            for row in ws.iter_rows():
                for cell in row:
                    if remaining <= 0:
                        break
                    if cell.data_type == "f" and not include_formulas:
                        continue
                    if not isinstance(cell.value, str):
                        continue
                    new_val, c = pat.subn(replace, cell.value)
                    if c:
                        cell.value = new_val
                        total_applied += min(c, remaining)
                        locations.append(
                            {"sheet": name, "cell": cell.coordinate, "count": c}
                        )
                        remaining -= c
        wb.save(str(p))
    elif ext == "pptx":
        from pptx import Presentation

        prs = Presentation(str(p))
        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                for p_idx, para in enumerate(shape.text_frame.paragraphs):
                    if remaining <= 0:
                        break
                    a, s = _replace_in_paragraph_runs(para, pat, replace)
                    if a:
                        total_applied += min(a, remaining)
                        total_skipped += s
                        locations.append(
                            {
                                "slide_index": s_idx,
                                "shape_index": sh_idx,
                                "paragraph_index": p_idx,
                                "count": a,
                            }
                        )
                        remaining -= a
            if include_notes and slide.has_notes_slide:
                for p_idx, para in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                    if remaining <= 0:
                        break
                    a, s = _replace_in_paragraph_runs(para, pat, replace)
                    if a:
                        total_applied += min(a, remaining)
                        total_skipped += s
                        locations.append(
                            {
                                "slide_index": s_idx,
                                "in_notes": True,
                                "paragraph_index": p_idx,
                                "count": a,
                            }
                        )
                        remaining -= a
        prs.save(str(p))
    else:
        raise ValueError(
            f"Unsupported format for replace: '.{ext}'.\n"
            f"Supported formats: docx, xlsx, pptx."
        )

    return {
        "replacements_made": total_applied,
        "skipped_multi_run": total_skipped,
        "locations": locations,
        "dry_run": False,
    }
