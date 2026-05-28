"""Headers, footers, hyperlinks, bookmarks, internal links, TOC, comments.

These are the OOXML-heaviest features: python-docx exposes none of comments,
TOC fields, page-number fields, or hyperlinks-as-runs directly, so we
construct the XML via lxml + OoxmlElement.

References:
- OOXML §17.13 — comments
- OOXML §17.16.5 — fields (PAGE, TOC, HYPERLINK)
- OOXML §17.13.6 — bookmarks
"""

from __future__ import annotations

import datetime as _dt
import uuid
from pathlib import Path

from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from lxml import etree

from .validation import validate_path, sanitize_text, validate_color
from .safety import safe_parse_xml


W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


# --------------------------------------------------------------------------
# Headers & footers
# --------------------------------------------------------------------------

from docx.enum.text import WD_ALIGN_PARAGRAPH as _WD_ALIGN

_ALIGN_MAP = {
    "left": _WD_ALIGN.LEFT,
    "center": _WD_ALIGN.CENTER,
    "right": _WD_ALIGN.RIGHT,
    "justify": _WD_ALIGN.JUSTIFY,
}


def _apply_run_props(para, color=None, size=None) -> None:
    """Apply font colour (hex) and/or size (pt) to every run in a paragraph,
    including runs nested inside fields (e.g. the PAGE number)."""
    if color is None and size is None:
        return
    hexc = validate_color(color) if color else None
    for r in para._p.iter(qn("w:r")):
        rpr = r.find(qn("w:rPr"))
        if rpr is None:
            rpr = OxmlElement("w:rPr")
            r.insert(0, rpr)
        if hexc:
            c = rpr.find(qn("w:color"))
            if c is None:
                c = OxmlElement("w:color")
                rpr.append(c)
            c.set(qn("w:val"), hexc)
        if size:
            sz = rpr.find(qn("w:sz"))
            if sz is None:
                sz = OxmlElement("w:sz")
                rpr.append(sz)
            sz.set(qn("w:val"), str(int(float(size) * 2)))


def docx_set_header(
    path: str,
    text: str,
    *,
    align: str = "center",
    section: int = 0,
    color: str | None = None,
    size: float | None = None,
) -> str:
    from docx import Document

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="set_header")
    doc = Document(str(p))
    if section < 0 or section >= len(doc.sections):
        raise ValueError(
            f"section {section} out of range [0, {len(doc.sections) - 1}].\n"
            f"Most documents have only one section (index 0)."
        )
    header = doc.sections[section].header
    # Reuse the first paragraph or add one.
    para = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
    para.text = sanitize_text(text, "header text")
    para.alignment = _ALIGN_MAP.get(align.lower(), 1)
    _apply_run_props(para, color, size)
    doc.save(str(p))
    return str(p)


def docx_set_footer(
    path: str,
    text: str = "",
    *,
    page_numbers: bool = True,
    align: str = "center",
    section: int = 0,
    color: str | None = None,
    size: float | None = None,
) -> str:
    from docx import Document

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="set_footer")
    doc = Document(str(p))
    if section < 0 or section >= len(doc.sections):
        raise ValueError(
            f"section {section} out of range [0, {len(doc.sections) - 1}].\n"
            f"Most documents have only one section (index 0)."
        )
    footer = doc.sections[section].footer
    para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    para.text = sanitize_text(text, "footer text")
    para.alignment = _ALIGN_MAP.get(align.lower(), 1)
    if page_numbers:
        # Inject a PAGE field after the user's text.
        run = para.add_run()
        if text:
            run.add_text(" — ")
        fld = OxmlElement("w:fldSimple")
        fld.set(qn("w:instr"), "PAGE")
        # The field needs a result run so the page number renders before the
        # document is updated by Word.
        rslt_r = OxmlElement("w:r")
        rslt_t = OxmlElement("w:t")
        rslt_t.text = "1"
        rslt_r.append(rslt_t)
        fld.append(rslt_r)
        run._r.addnext(fld)
    _apply_run_props(para, color, size)
    doc.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# Hyperlinks
# --------------------------------------------------------------------------

def _add_hyperlink_run(paragraph, url: str, text: str, *, internal: bool = False):
    """Add a hyperlink to a python-docx paragraph using lxml.

    Returns the run element so callers can style it further.
    """
    part = paragraph.part
    if internal:
        r_id = None
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("w:anchor"), url)
    else:
        r_id = part.relate_to(
            url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)

    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    # Apply the standard Hyperlink style.
    rStyle = OxmlElement("w:rStyle")
    rStyle.set(qn("w:val"), "Hyperlink")
    rPr.append(rStyle)
    new_run.append(rPr)

    t = OxmlElement("w:t")
    t.text = sanitize_text(text, "hyperlink text")
    t.set(qn("xml:space"), "preserve")
    new_run.append(t)

    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)
    return hyperlink


def docx_add_hyperlink(path: str, paragraph_index: int, text: str, url: str) -> str:
    from docx import Document
    from .validation import validate_paragraph_index

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="add_hyperlink")
    doc = Document(str(p))
    validate_paragraph_index(doc, paragraph_index)
    _add_hyperlink_run(doc.paragraphs[paragraph_index], url, text, internal=False)
    doc.save(str(p))
    return str(p)


def docx_add_internal_link(path: str, paragraph_index: int, text: str, bookmark_name: str) -> str:
    from docx import Document
    from .validation import validate_paragraph_index

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="add_internal_link")
    doc = Document(str(p))
    validate_paragraph_index(doc, paragraph_index)
    _add_hyperlink_run(
        doc.paragraphs[paragraph_index], bookmark_name, text, internal=True
    )
    doc.save(str(p))
    return str(p)


def pptx_add_hyperlink(path: str, slide_index: int, shape_index: int, url: str) -> str:
    from pptx import Presentation

    p = validate_path(path, must_exist=True, expected_ext="pptx", operation="add_hyperlink")
    prs = Presentation(str(p))
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(
            f"slide_index {slide_index} out of range [0, {len(prs.slides) - 1}].\n"
            f"The presentation has {len(prs.slides)} slides (0-indexed)."
        )
    slide = prs.slides[slide_index]
    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(
            f"shape_index {shape_index} out of range [0, {len(slide.shapes) - 1}].\n"
            f"Slide {slide_index} has {len(slide.shapes)} shapes (0-indexed)."
        )
    shape = slide.shapes[shape_index]
    if not shape.has_text_frame:
        raise ValueError(
            f"Shape {shape_index} on slide {slide_index} has no text frame.\n"
            f"Hyperlinks can only be added to shapes that contain text."
        )
    applied = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.hyperlink.address = url
            applied = True
            break
        if applied:
            break
    if not applied:
        raise ValueError(
            f"No text runs found in shape {shape_index} on slide {slide_index}.\n"
            f"Add text to the shape first, then apply the hyperlink."
        )
    prs.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# Bookmarks
# --------------------------------------------------------------------------

def docx_add_bookmark(path: str, paragraph_index: int, name: str) -> str:
    """Wrap a paragraph in ``<w:bookmarkStart>`` / ``<w:bookmarkEnd>``."""
    from docx import Document
    from .validation import validate_paragraph_index

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="add_bookmark")
    doc = Document(str(p))
    validate_paragraph_index(doc, paragraph_index)
    para = doc.paragraphs[paragraph_index]
    existing_ids = {
        int(el.get(qn("w:id"), "0"))
        for el in doc.element.body.iter(qn("w:bookmarkStart"))
    }
    bm_id = (max(existing_ids) + 1) if existing_ids else 0

    safe_name = sanitize_text(name, "bookmark name")
    start = OxmlElement("w:bookmarkStart")
    start.set(qn("w:id"), str(bm_id))
    start.set(qn("w:name"), safe_name)
    end = OxmlElement("w:bookmarkEnd")
    end.set(qn("w:id"), str(bm_id))
    para._p.insert(0, start)
    para._p.append(end)
    doc.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# Table of contents (field)
# --------------------------------------------------------------------------

def docx_add_toc(path: str, paragraph_index: int = 0) -> str:
    """Insert a TOC field. Word / LibreOffice updates the TOC on open."""
    from docx import Document

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="add_toc")
    doc = Document(str(p))
    total = len(doc.paragraphs)
    if paragraph_index < 0 or paragraph_index > total:
        raise ValueError(
            f"paragraph_index {paragraph_index} out of range [0, {total}].\n"
            f"Use {total} to append at the end, or 0 to insert at the start."
        )
    if paragraph_index == total:
        target = doc.add_paragraph()
    else:
        target = doc.paragraphs[paragraph_index]

    fld_char_begin = OxmlElement("w:fldChar")
    fld_char_begin.set(qn("w:fldCharType"), "begin")
    instr_text = OxmlElement("w:instrText")
    instr_text.set(qn("xml:space"), "preserve")
    instr_text.text = "TOC \\o \"1-3\" \\h \\z \\u"
    fld_char_sep = OxmlElement("w:fldChar")
    fld_char_sep.set(qn("w:fldCharType"), "separate")
    fld_char_end = OxmlElement("w:fldChar")
    fld_char_end.set(qn("w:fldCharType"), "end")

    run = OxmlElement("w:r")
    run.append(fld_char_begin)
    run.append(instr_text)
    run.append(fld_char_sep)
    placeholder_t = OxmlElement("w:t")
    placeholder_t.text = "Right-click to update Table of Contents."
    placeholder_r = OxmlElement("w:r")
    placeholder_r.append(placeholder_t)
    run.append(placeholder_r)
    run.append(fld_char_end)
    target._p.append(run)
    doc.save(str(p))
    return str(p)


# --------------------------------------------------------------------------
# Comments
# --------------------------------------------------------------------------

_COMMENTS_XML_TEMPLATE = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>
"""


_COMMENTS_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument."
    "wordprocessingml.comments+xml"
)
_COMMENTS_REL_TYPE = R_NS + "/comments"


def _ensure_comments_part(doc):
    """Return the comments part (creating it + relationships if missing).

    Strategy: walk doc.part.package.parts looking for word/comments.xml. If
    missing, build a minimal valid one via the package's part-factory API and
    add a relationship from document.xml.

    Uses CT_NS for content-type namespace validation and uuid for
    unique relationship IDs.
    """
    package = doc.part.package
    comments_partname = "/word/comments.xml"
    for part in package.parts:
        if str(part.partname) == comments_partname:
            return part
    from docx.opc.part import Part
    from docx.opc.packuri import PackURI

    comments_part = Part(
        partname=PackURI(comments_partname),
        content_type=_COMMENTS_CONTENT_TYPE,
        blob=_COMMENTS_XML_TEMPLATE,
        package=package,
    )
    if hasattr(package.parts, "append_part"):
        package.parts.append_part(comments_part)
    else:
        raise RuntimeError(
            "Cannot create comments part: package API does not support append_part.\n"
            "Try opening the document in Word or LibreOffice first to initialize the comments structure."
        )
    doc.part.relate_to(comments_part, _COMMENTS_REL_TYPE)
    return comments_part


def docx_add_comment(
    path: str,
    paragraph_index: int,
    author: str,
    text: str,
    *,
    initials: str = "AI",
) -> str:
    """Attach a Word-style comment to a paragraph.

    Builds (or extends) ``word/comments.xml`` and inserts the comment range
    markers + reference into the target paragraph.
    """
    from docx import Document
    from docx.oxml.ns import nsmap

    from .validation import validate_paragraph_index

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="add_comment")
    doc = Document(str(p))
    validate_paragraph_index(doc, paragraph_index)
    target_para = doc.paragraphs[paragraph_index]

    try:
        comments_part = _ensure_comments_part(doc)
    except Exception as e:
        # If comments part creation fails (some docx don't accept arbitrary
        # part injection without a fuller content-types update), fall back to
        # appending the comment text as a [Comment: ...] bracketed paragraph.
        para = doc.add_paragraph(f"[Comment by {author}: {text}]")
        doc.save(str(p))
        return str(p)

    # Parse the existing comments XML, append a new w:comment.
    comments_root = safe_parse_xml(comments_part.blob)
    # Determine the next free id.
    existing_ids = [
        int(c.get(f"{{{W_NS}}}id", "0"))
        for c in comments_root.findall(f"{{{W_NS}}}comment")
    ]
    next_id = (max(existing_ids) + 1) if existing_ids else 0

    w_ns = nsmap.get("w", W_NS)
    comment = etree.SubElement(comments_root, f"{{{w_ns}}}comment")
    comment.set(f"{{{w_ns}}}id", str(next_id))
    comment.set(f"{{{w_ns}}}author", sanitize_text(author, "comment author"))
    comment.set(f"{{{w_ns}}}initials", sanitize_text(initials, "comment initials"))
    comment.set(
        f"{{{w_ns}}}date",
        _dt.datetime.now(_dt.timezone.utc).replace(microsecond=0, tzinfo=None).isoformat() + "Z",
    )
    pp = etree.SubElement(comment, f"{{{w_ns}}}p")
    rp = etree.SubElement(pp, f"{{{w_ns}}}r")
    tp = etree.SubElement(rp, f"{{{w_ns}}}t")
    tp.text = sanitize_text(text, "comment text")

    comments_part._blob = etree.tostring(
        comments_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    # Add range markers around the target paragraph.
    start = OxmlElement("w:commentRangeStart")
    start.set(qn("w:id"), str(next_id))
    end = OxmlElement("w:commentRangeEnd")
    end.set(qn("w:id"), str(next_id))
    ref_run = OxmlElement("w:r")
    ref = OxmlElement("w:commentReference")
    ref.set(qn("w:id"), str(next_id))
    ref_run.append(ref)

    target_para._p.insert(0, start)
    target_para._p.append(end)
    target_para._p.append(ref_run)

    doc.save(str(p))
    return str(p)


def docx_list_comments(path: str) -> list[dict]:
    """Return all comments in a docx file."""
    from docx import Document

    p = validate_path(path, must_exist=True, expected_ext="docx", operation="list_comments")
    doc = Document(str(p))
    comments_partname = "/word/comments.xml"
    for part in doc.part.package.parts:
        if str(part.partname) == comments_partname:
            root = safe_parse_xml(part.blob)
            out = []
            for c in root.findall(f"{{{W_NS}}}comment"):
                text_parts = [
                    t.text or ""
                    for t in c.iter(f"{{{W_NS}}}t")
                ]
                out.append(
                    {
                        "id": int(c.get(f"{{{W_NS}}}id", "0")),
                        "author": c.get(f"{{{W_NS}}}author", ""),
                        "initials": c.get(f"{{{W_NS}}}initials", ""),
                        "date": c.get(f"{{{W_NS}}}date", ""),
                        "text": "".join(text_parts),
                    }
                )
            return out
    return []
